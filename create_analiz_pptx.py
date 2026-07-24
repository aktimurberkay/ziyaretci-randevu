#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Ziyaretçi ve Randevu Yönetim Sistemi - Analiz Dokümanı (Akış Şemalı Versiyon)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ─── RENK PALETİ ───
KOYU_MAVI    = RGBColor(0x1B, 0x3A, 0x5C)
ORTA_MAVI    = RGBColor(0x2C, 0x5F, 0x8A)
ACIK_MAVI    = RGBColor(0x3A, 0x7C, 0xBD)
BEYAZ        = RGBColor(0xFF, 0xFF, 0xFF)
SIYAH        = RGBColor(0x00, 0x00, 0x00)
KOYU_GRI     = RGBColor(0x33, 0x33, 0x33)
ORTA_GRI     = RGBColor(0x66, 0x66, 0x66)
ACIK_GRI     = RGBColor(0xF2, 0xF2, 0xF2)
TABLO_BASLIK = RGBColor(0x1B, 0x3A, 0x5C)
ACIK_SARI    = RGBColor(0xFF, 0xF3, 0xCD)
ACIK_MAVI_BG = RGBColor(0xD6, 0xEC, 0xF8)
ACIK_YESIL   = RGBColor(0xD4, 0xED, 0xDA)
ACIK_KIRMIZI = RGBColor(0xF8, 0xD7, 0xDA)

A4_WIDTH  = Cm(21.0)
A4_HEIGHT = Cm(29.7)

prs = Presentation()
prs.slide_width = A4_WIDTH
prs.slide_height = A4_HEIGHT

# ═══════════════════════════════════════
# GENEL YARDIMCI FONKSİYONLAR
# ═══════════════════════════════════════
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                color=SIYAH, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_footer(slide, page_num):
    add_textbox(slide, Cm(1), Cm(28.0), Cm(6), Cm(0.8),
                "Dahili Kullanım", font_size=8, color=ORTA_GRI)
    add_textbox(slide, Cm(17), Cm(28.0), Cm(3), Cm(0.8),
                str(page_num), font_size=8, color=ORTA_GRI, alignment=PP_ALIGN.RIGHT)

def add_blue_header_bar(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), A4_WIDTH, Cm(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KOYU_MAVI
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BEYAZ
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Cm(0.4)

def add_section_title(slide, text, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), top, Cm(19), Cm(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORTA_MAVI
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BEYAZ
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(0.3)

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = 'Calibri'
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BEYAZ
                else:
                    paragraph.font.color.rgb = KOYU_GRI
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLO_BASLIK
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACIK_GRI
    return table_shape


# ═══════════════════════════════════════
# AKIŞ ŞEMASI YARDIMCI FONKSİYONLAR
# ═══════════════════════════════════════
def fc_shape(slide, shape_type, left, top, width, height, text,
             fill_color=BEYAZ, text_color=SIYAH, border_color=KOYU_MAVI,
             font_size=8, bold=False):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return shape

def fc_terminator(slide, cx, top, text, w=Cm(3.8), h=Cm(1.0)):
    return fc_shape(slide, MSO_SHAPE.FLOWCHART_TERMINATOR,
                    cx - w//2, top, w, h, text,
                    fill_color=KOYU_MAVI, text_color=BEYAZ, font_size=9, bold=True)

def fc_process(slide, cx, top, text, w=Cm(5.5), h=Cm(1.3)):
    return fc_shape(slide, MSO_SHAPE.FLOWCHART_PROCESS,
                    cx - w//2, top, w, h, text,
                    fill_color=BEYAZ, text_color=KOYU_GRI, font_size=8, bold=True)

def fc_decision(slide, cx, top, text, w=Cm(5.0), h=Cm(2.4)):
    return fc_shape(slide, MSO_SHAPE.FLOWCHART_DECISION,
                    cx - w//2, top, w, h, text,
                    fill_color=ACIK_SARI, text_color=KOYU_GRI, font_size=7, bold=True)

def fc_data(slide, cx, top, text, w=Cm(6.0), h=Cm(1.4)):
    return fc_shape(slide, MSO_SHAPE.FLOWCHART_DATA,
                    cx - w//2, top, w, h, text,
                    fill_color=ACIK_MAVI_BG, text_color=KOYU_GRI, font_size=8, bold=True)

def fc_result(slide, left, top, text, fill_color, w=Cm(4.5), h=Cm(1.3)):
    return fc_shape(slide, MSO_SHAPE.FLOWCHART_PROCESS,
                    left, top, w, h, text,
                    fill_color=fill_color, text_color=KOYU_GRI, font_size=7, bold=True)

def arrow_down(slide, cx, start_y, end_y, color=SIYAH):
    t = Cm(0.05)
    a_w = Cm(0.35)
    a_h = Cm(0.25)
    line_len = end_y - start_y - a_h
    if line_len > 0:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - t//2, start_y, t, line_len)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  cx - a_w//2, end_y - a_h, a_w, a_h)
    tri.rotation = 180.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def arrow_right(slide, start_x, cy, end_x, color=SIYAH):
    t = Cm(0.05)
    a_w = Cm(0.25)
    a_h = Cm(0.35)
    line_len = end_x - start_x - a_w
    if line_len > 0:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, cy - t//2, line_len, t)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  end_x - a_w, cy - a_h//2, a_w, a_h)
    tri.rotation = 90.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def arrow_left(slide, start_x, cy, end_x, color=SIYAH):
    t = Cm(0.05)
    a_w = Cm(0.25)
    a_h = Cm(0.35)
    line_len = start_x - end_x - a_w
    if line_len > 0:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, end_x + a_w, cy - t//2, line_len, t)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  end_x, cy - a_h//2, a_w, a_h)
    tri.rotation = 270.0
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def fc_label(slide, left, top, text, font_size=7):
    add_textbox(slide, left, top, Cm(2.5), Cm(0.5), text,
                font_size=font_size, bold=True, color=ORTA_GRI)


# ═══════════════════════════════════════
# SLAYT 1: KAPAK
# ═══════════════════════════════════════
slide = add_blank_slide()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), A4_WIDTH, Cm(3))
shape.fill.solid(); shape.fill.fore_color.rgb = KOYU_MAVI; shape.line.fill.background()
shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(25), A4_WIDTH, Cm(4.7))
shape2.fill.solid(); shape2.fill.fore_color.rgb = KOYU_MAVI; shape2.line.fill.background()
add_textbox(slide, Cm(2), Cm(8), Cm(17), Cm(3),
            "Ziyaretçi ve Randevu\nYönetim Sistemi",
            font_size=32, bold=True, color=KOYU_MAVI, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(12.5), Cm(17), Cm(1.5),
            "Analiz Dokümanı",
            font_size=22, bold=False, color=ORTA_MAVI, alignment=PP_ALIGN.CENTER)
shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(7), Cm(14.5), Cm(7), Cm(0.08))
shape3.fill.solid(); shape3.fill.fore_color.rgb = ACIK_MAVI; shape3.line.fill.background()
add_textbox(slide, Cm(2), Cm(16), Cm(17), Cm(1),
            "Bilişim Teknolojileri Bölümü", font_size=14, color=ORTA_GRI, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(17.5), Cm(17), Cm(1),
            "Stajyer", font_size=12, color=ORTA_GRI, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(26), Cm(17), Cm(1),
            "Ziyaretçi Kayıt Sistemi Analiz Dokümanı",
            font_size=11, color=BEYAZ, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(27.5), Cm(17), Cm(1),
            "Dahili Kullanım", font_size=9, color=RGBColor(0x88,0xAA,0xCC), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLAYT 2: GENEL PROJE BİLGİLERİ
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "GENEL PROJE BİLGİLERİ")
add_section_title(slide, "Proje Bilgileri", Cm(3))
add_table(slide, Cm(1), Cm(4.2), Cm(19), Cm(4), 5, 4,
    [["Alan", "Değer", "Alan", "Değer"],
     ["Proje Tipi", "Ziyaretçi ve Randevu\nYönetim Sistemi Analizi", "Ürün", "Web Uygulaması"],
     ["Proje No", "1", "Son Revize Tarihi", "19.07.2026"],
     ["Hazırlayan", "Berkay Aktimur", "Öncelik", "Normal"],
     ["Hazırlama Tarihi", "19.07.2026", "Versiyon", "1.0"]],
    col_widths=[Cm(4), Cm(6), Cm(4), Cm(5)])
add_section_title(slide, "Onay Bilgileri", Cm(9.5))
add_table(slide, Cm(1), Cm(10.7), Cm(19), Cm(3), 4, 4,
    [["Onaylayan", "Teslim Tarihi", "Onay Tarihi", "Açıklama"],
     ["", "", "", ""], ["", "", "", ""], ["", "", "", ""]],
    col_widths=[Cm(5), Cm(4.5), Cm(4.5), Cm(5)])
add_footer(slide, 2)


# ═══════════════════════════════════════
# SLAYT 3: İÇİNDEKİLER
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "İÇİNDEKİLER TABLOSU")
icerik_items = [
    ("Genel Proje Bilgileri", "2"),
    ("Talep ve İhtiyaç / Amaç ve Kapsam", "4"),
    ("Sistem Mimarisi ve Teknik Altyapı", "5"),
    ("Veritabanı Yapısı", "6"),
    ("İş Akış Şeması – Giriş ve Kayıt Süreci", "7"),
    ("İş Akış Şeması – Onay Süreci", "8"),
    ("İş Akış Şeması – Silme Süreci", "9"),
    ("İş Akış Şeması – Çıkış Süreci", "10"),
    ("Kullanıcı Rolleri ve Yetkilendirme", "11"),
    ("Sayfa Ön Tasarımları – İş Kuralları (1)", "12"),
    ("Sayfa Ön Tasarımları – İş Kuralları (2)", "13"),
    ("Raporlama", "14"),
]
txBox = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(17), Cm(22))
tf = txBox.text_frame; tf.word_wrap = True
for i, (baslik, sayfa) in enumerate(icerik_items):
    dots = "." * (75 - len(baslik) - len(sayfa))
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{baslik}  {dots}  {sayfa}"
    p.font.size = Pt(11); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(10)
add_footer(slide, 3)


# ═══════════════════════════════════════
# SLAYT 4: TALEP / AMAÇ / KAPSAM
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "TALEP VE İHTİYAÇ  /  AMAÇ VE KAPSAM")
add_section_title(slide, "Talep ve İhtiyaç", Cm(3))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(4.2), Cm(18), Cm(3.5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Kurum bünyesindeki lokasyonlarda, ziyaretçi kayıtlarının düzenli bir şekilde oluşturulması, "
    "yönetilmesi, onaylanması ve arşivlenmesi ihtiyacı üzerine bu proje geliştirilmiştir. "
    "Ziyaretçi kayıtlarının web üzerinden kolayca kaydedilmesi, randevu ve onay süreçlerinin "
    "takip edilmesi, personeller arası iletişimin sağlanması ve geçmiş kayıtların "
    "raporlanabilir şekilde saklanması hedeflenmektedir.")
tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = KOYU_GRI

add_section_title(slide, "Amaç", Cm(8))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(9.2), Cm(18), Cm(4))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Tüm kurumsal lokasyonlarda ziyaretçi kayıtlarının sistematik bir şekilde oluşturulması, "
    "yönetilmesi ve takip edilmesi amaçlanmaktadır. Sistem, ziyaretçi bilgilerinin (T.C. Kimlik, ad, "
    "soyad, iletişim, ziyaret nedeni, görüşülecek kişi) kaydedilmesini sağlayan bir kayıt formu "
    "sunacaktır. Kayıtların yetkili kullanıcılar (admin/resepsiyon) veya sınırlı yetkili kullanıcılar "
    "(personel/sekreter) tarafından doldurulması, onaylanması veya reddedilmesi mümkün olacaktır.")
tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = KOYU_GRI

add_section_title(slide, "Kapsam", Cm(13.5))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(14.7), Cm(18), Cm(12))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Sistem, web tabanlı mimaride çalışmakta olup mobil ve masaüstü uyumludur:"
tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = KOYU_GRI
kapsam = [
    "Randevular, admin/resepsiyon tarafından doğrudan oluşturulabilir; personel/sekreter 'onay bekleyen' statüsünde kayıt oluşturur.",
    "Onay bekleyen kayıtlar, ilgili personele bildirim ve/veya e-posta ile iletilir.",
    "Onaylanan randevular takvim sistemine otomatik yansıtılır.",
    "Geçmiş kayıtlar ve silinen kayıtlar ayrı bölümlerde arşivlenir.",
    "CSV ve PDF formatlarında dışa aktarma desteği bulunur.",
    "Kara liste (blacklist) kontrolü ile güvenlik riski oluşturan kişilerin girişi engellenir.",
    "Denetim izi (Audit Log) ile tüm işlemler kayıt altına alınır.",
    "Personeller arası mesajlaşma modülü ile kurum içi iletişim sağlanır.",
]
for item in kapsam:
    p = tf.add_paragraph(); p.text = f"• {item}"
    p.font.size = Pt(9); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(3)
add_footer(slide, 4)


# ═══════════════════════════════════════
# SLAYT 5: SİSTEM MİMARİSİ
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "SİSTEM MİMARİSİ VE TEKNİK ALTYAPI")
add_section_title(slide, "Backend (Sunucu Tarafı)", Cm(3))
add_table(slide, Cm(1), Cm(4.2), Cm(19), Cm(7), 9, 3,
    [["Bileşen", "Teknoloji", "Açıklama"],
     ["Programlama Dili", "Python 3.x", "Ana uygulama dili"],
     ["Web Framework", "Flask", "Hafif ve esnek web çatısı"],
     ["ORM", "SQLAlchemy (Flask-SQLAlchemy)", "Veritabanı nesne-ilişki yönetimi"],
     ["Kimlik Doğrulama", "Flask-Login", "Oturum yönetimi ve kullanıcı doğrulama"],
     ["Şifreleme", "Werkzeug Security", "Parola hash'leme (pbkdf2)"],
     ["E-posta", "Flask-Mail (SMTP)", "Bildirim ve randevu e-postaları"],
     ["Güvenlik", "Flask-Limiter", "Rate-limiting / DDoS koruması"],
     ["Raporlama", "ReportLab, CSV modülü", "PDF ve CSV dışa aktarma"]],
    col_widths=[Cm(5), Cm(6.5), Cm(7.5)])
add_section_title(slide, "Frontend (İstemci Tarafı)", Cm(12.5))
add_table(slide, Cm(1), Cm(13.7), Cm(19), Cm(4), 5, 3,
    [["Bileşen", "Teknoloji", "Açıklama"],
     ["Şablon Motoru", "Jinja2 (Flask entegre)", "Sunucu taraflı HTML oluşturma"],
     ["Yapı", "HTML5", "Semantik ve erişilebilir işaretleme"],
     ["Stil", "CSS3 (Responsive)", "Mobil ve masaüstü uyumlu tasarım"],
     ["Etkileşim", "JavaScript (Vanilla)", "Dinamik arayüz etkileşimleri"]],
    col_widths=[Cm(5), Cm(6.5), Cm(7.5)])
add_section_title(slide, "Veritabanı", Cm(18.5))
add_table(slide, Cm(1), Cm(19.7), Cm(19), Cm(3.5), 4, 3,
    [["Bileşen", "Teknoloji", "Açıklama"],
     ["RDBMS", "SQLite", "Geliştirme ve operasyonel kullanım"],
     ["ORM Katmanı", "SQLAlchemy", "İlişkisel veritabanı soyutlama"],
     ["Zaman Dilimi", "pytz (Europe/Istanbul)", "Türkiye saatine uyumlu kayıtlar"]],
    col_widths=[Cm(5), Cm(6.5), Cm(7.5)])
add_footer(slide, 5)


# ═══════════════════════════════════════
# SLAYT 6: VERİTABANI YAPISI
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "VERİTABANI YAPISI")
add_section_title(slide, "Temel Veri Tabloları (Modeller)", Cm(3))
add_table(slide, Cm(0.5), Cm(4.2), Cm(20), Cm(15), 10, 3,
    [["Tablo (Model)", "Açıklama", "Temel Alanlar"],
     ["Kullanici", "Sistem kullanıcıları\n(Admin, Personel, Resepsiyon, Sekreter)", "kullanici_adi, ad_soyad, eposta,\nrol, departman, aktif_durum"],
     ["Ziyaretci", "Fiziksel ziyaretçi bilgileri", "ad_soyad, tc_kimlik, eposta,\nsirket, telefon, kvkk_onayi"],
     ["Randevu", "Ziyaretçi-Personel arası\nrandevu kaydı ve durumları", "ziyaretci_id, personel_id, tarih_saat,\ndurum, giris_saati, cikis_saati, notlar"],
     ["AuditLog", "Sistem denetim izi\n(tüm işlem logları)", "kullanici_id, islem_tipi, hedef_tablo,\nhedef_id, detay, tarih"],
     ["TakvimEtkinlik", "Personel takvim aktiviteleri", "sahibi_id, baslik, baslangic,\nbitis, tip, aciklama"],
     ["Mesaj", "Sistem içi mesajlaşma", "gonderen_id, alici_id, icerik,\nokundu, tarih_saat"],
     ["KaraListe", "Girişi yasaklı kişiler", "tc_kimlik, eposta, sebep,\neklenme_tarihi, ekleyen_id"],
     ["Bildirim", "Kullanıcı bildirimleri", "kullanici_id, icerik, okundu,\ntarih_saat"],
     ["ZiyaretciMailGecmisi", "Ziyaretçilere gönderilen\ne-posta logları", "ziyaretci_id, gonderen_kullanici_id,\nkonu, icerik"]],
    col_widths=[Cm(5), Cm(6.5), Cm(8.5)])

txBox = slide.shapes.add_textbox(Cm(1), Cm(20.5), Cm(19), Cm(5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Tablo İlişkileri:"
tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI
for rel in [
    "• Kullanici (1) ── (N) Randevu  (personel_id, olusturan_id)",
    "• Ziyaretci (1) ── (N) Randevu  (ziyaretci_id)",
    "• Kullanici (1) ── (N) TakvimEtkinlik  (sahibi_id, olusturan_id)",
    "• Kullanici (1) ── (N) Mesaj  (gonderen_id, alici_id)",
    "• Kullanici (1) ── (N) AuditLog / KaraListe / Bildirim",
    "• Kullanici (1) ── (N) Kullanici  (bagli_yonetici_id → Sekreter-Yönetici)"]:
    p = tf.add_paragraph(); p.text = rel
    p.font.size = Pt(8.5); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(2)
add_footer(slide, 6)


# ═══════════════════════════════════════════════════════════════
# SLAYT 7: İŞ AKIŞ ŞEMASI – GİRİŞ VE KAYIT SÜRECİ (ŞEKİLLİ)
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "İŞ AKIŞ ŞEMALARI")
add_section_title(slide, "Giriş ve Kayıt Süreci Akış Şeması", Cm(2.8))

CX = Cm(10.5)  # Sayfa yatay merkezi
SW = Cm(5.5)   # Standart şekil genişliği
DW = Cm(5.0)   # Decision genişliği
DH = Cm(2.2)   # Decision yüksekliği

# 1) BAŞLAT
y = Cm(4.2)
fc_terminator(slide, CX, y, "BAŞLAT")
arrow_down(slide, CX, y + Cm(1.0), y + Cm(1.5))

# 2) GİRİŞ EKRANI
y = Cm(5.7)
fc_data(slide, CX, y, "GİRİŞ EKRANI\n(KULLANICI ADI & ŞİFRE)", w=Cm(6.5))
arrow_down(slide, CX, y + Cm(1.4), y + Cm(1.9))

# 3) BİLGİLER DOĞRU MU?
y_dec1 = Cm(7.6)
fc_decision(slide, CX, y_dec1, "BİLGİLER\nDOĞRU MU?", w=DW, h=DH)

# Sağ dal → HATALI GİRİŞ
dec1_right = CX + DW//2
dec1_cy = y_dec1 + DH//2
arrow_right(slide, dec1_right, dec1_cy, Cm(15.5))
fc_result(slide, Cm(15.5), dec1_cy - Cm(0.65), "HATALI GİRİŞ\n(HATA MESAJI)", ACIK_KIRMIZI, w=Cm(4.5))
fc_label(slide, dec1_right + Cm(0.2), dec1_cy - Cm(0.7), "YANLIŞ")

# Aşağı dal → DOĞRU
fc_label(slide, CX + Cm(0.3), y_dec1 + DH - Cm(0.1), "DOĞRU")
arrow_down(slide, CX, y_dec1 + DH, y_dec1 + DH + Cm(0.5))

# 4) KULLANICI TİPİNE GÖRE ARAYÜZ
y = y_dec1 + DH + Cm(0.5)
fc_process(slide, CX, y, "KULLANICI TİPİNE\nGÖRE ARAYÜZ")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 5) ZİYARETÇİ BİLGİLERİ GİRİLİR
y = y + Cm(1.8)
fc_process(slide, CX, y, "ZİYARETÇİ BİLGİLERİ\nGİRİLİR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 6) KULLANICI TİPİ ADMİN Mİ USER MI?
y_dec2 = y + Cm(1.8)
fc_decision(slide, CX, y_dec2, "KULLANICI TİPİ\nADMİN Mİ\nUSER MI?", w=DW, h=DH)

# Sol dal → ADMIN
dec2_left = CX - DW//2
dec2_cy = y_dec2 + DH//2
arrow_left(slide, dec2_left, dec2_cy, Cm(1.0) + Cm(4.5))
fc_result(slide, Cm(1.0), dec2_cy - Cm(0.65), "ADMIN:\nKAYIT → RECORDS", ACIK_YESIL, w=Cm(4.5))
fc_label(slide, dec2_left - Cm(2.8), dec2_cy - Cm(0.7), "ADMİN")

# Sağ dal → USER
dec2_right = CX + DW//2
arrow_right(slide, dec2_right, dec2_cy, Cm(15.5))
fc_result(slide, Cm(15.5), dec2_cy - Cm(0.65), "USER:\nKAYIT → PENDING", ACIK_SARI, w=Cm(4.5))
fc_label(slide, dec2_right + Cm(0.2), dec2_cy - Cm(0.7), "USER")

# Aşağı ok
arrow_down(slide, CX, y_dec2 + DH, y_dec2 + DH + Cm(0.5))

# 7) E-POSTA / BİLDİRİM
y = y_dec2 + DH + Cm(0.5)
fc_process(slide, CX, y, "E-POSTA / BİLDİRİM\nGÖNDERİMİ")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 8) BİTİR
y = y + Cm(1.8)
fc_terminator(slide, CX, y, "BİTİR")

# Açıklama metni (alt kısım)
desc_y = y + Cm(1.5)
txBox = slide.shapes.add_textbox(Cm(1), desc_y, Cm(19), Cm(3.5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Kayıt süreci, kullanıcının sisteme giriş yapması ile başlar. Bilgiler doğrulanır, "
    "hatalıysa brute-force koruması devreye girer. Doğruysa rol bazlı arayüz gösterilir. "
    "Ziyaretçi bilgileri girilir ve kullanıcı tipi kontrol edilir. Admin doğrudan kaydederken, "
    "User kayıtları bekleyen listeye gönderilir. Her iki durumda da bildirim iletilir.")
tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = ORTA_GRI

add_footer(slide, 7)


# ═══════════════════════════════════════════════════════════════
# SLAYT 8: İŞ AKIŞ ŞEMASI – ONAY SÜRECİ (ŞEKİLLİ)
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "İŞ AKIŞ ŞEMALARI")
add_section_title(slide, "Onay Süreci Akış Şeması", Cm(2.8))

CX = Cm(10.5)

# 1) BAŞLAT
y = Cm(4.2)
fc_terminator(slide, CX, y, "BAŞLAT")
arrow_down(slide, CX, y + Cm(1.0), y + Cm(1.5))

# 2) PERSONEL GİRİŞ YAPAR
y = Cm(5.7)
fc_process(slide, CX, y, "PERSONEL SİSTEME\nGİRİŞ YAPAR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 3) BEKLEYEN RANDEVULARI GÖRÜNTÜLER
y = Cm(7.5)
fc_process(slide, CX, y, "BEKLEYEN RANDEVULARI\nGÖRÜNTÜLER")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 4) ZİYARETÇİ BİLGİLERİNİ İNCELER
y = Cm(9.3)
fc_process(slide, CX, y, "ZİYARETÇİ BİLGİLERİNİ\nİNCELER")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 5) ONAYLA MI REDDET Mİ?
y_dec = Cm(11.1)
DH_onay = Cm(2.5)
fc_decision(slide, CX, y_dec, "ONAYLA MI\nREDDET Mİ?", w=DW, h=DH_onay)

# Sol dal → ONAYLA
dec_left = CX - DW//2
dec_cy = y_dec + DH_onay//2
arrow_left(slide, dec_left, dec_cy, Cm(1.0) + Cm(4.5))
fc_result(slide, Cm(1.0), dec_cy - Cm(0.65), "RANDEVU ONAYLANDI\nTAKVİME EKLENDİ", ACIK_YESIL, w=Cm(4.5))
fc_label(slide, dec_left - Cm(2.8), dec_cy - Cm(0.7), "ONAYLA")

# Sağ dal → REDDET
dec_right = CX + DW//2
arrow_right(slide, dec_right, dec_cy, Cm(15.5))
fc_result(slide, Cm(15.5), dec_cy - Cm(0.65), "RED AÇIKLAMASI GİRİLDİ\nRANDEVU REDDEDİLDİ", ACIK_KIRMIZI, w=Cm(4.5))
fc_label(slide, dec_right + Cm(0.2), dec_cy - Cm(0.7), "REDDET")

# Aşağı ok
arrow_down(slide, CX, y_dec + DH_onay, y_dec + DH_onay + Cm(0.5))

# 6) AUDIT LOG
y = y_dec + DH_onay + Cm(0.5)
fc_process(slide, CX, y, "AUDIT LOG KAYDI\nOLUŞTURULDU")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 7) BİLDİRİM GÖNDERİMİ
y = y + Cm(1.8)
fc_process(slide, CX, y, "BİLDİRİM / E-POSTA\nGÖNDERİLDİ")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.8))

# 8) BİTİR
y = y + Cm(1.8)
fc_terminator(slide, CX, y, "BİTİR")

# Açıklama
desc_y = y + Cm(1.5)
txBox = slide.shapes.add_textbox(Cm(1), desc_y, Cm(19), Cm(3.5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Personel, sisteme giriş yaparak kendisine atanmış bekleyen randevuları görüntüler. "
    "Ziyaretçi bilgilerini ve ziyaret nedenini inceledikten sonra onaylama veya reddetme kararı verir. "
    "Red durumunda bir açıklama girilmesi zorunludur. Onay durumunda randevu takvime otomatik eklenir. "
    "Tüm onay/red işlemleri Audit Log ile kayıt altına alınır ve ilgili taraflara bildirim gönderilir.")
tf.paragraphs[0].font.size = Pt(8); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = ORTA_GRI

add_footer(slide, 8)


# ═══════════════════════════════════════════════════════════════
# SLAYT 9: İŞ AKIŞ ŞEMASI – SİLME SÜRECİ (ŞEKİLLİ)
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "İŞ AKIŞ ŞEMALARI")
add_section_title(slide, "Silme Süreci Akış Şeması", Cm(2.8))

CX = Cm(10.5)

# 1) BAŞLAT
y = Cm(4.5)
fc_terminator(slide, CX, y, "BAŞLAT")
arrow_down(slide, CX, y + Cm(1.0), y + Cm(1.6))

# 2) ADMİN SİL BUTONUNA TIKLAR
y = Cm(6.1)
fc_process(slide, CX, y, "ADMİN \"SİL\" BUTONUNA\nTIKLAR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 3) SİSTEM GEREKÇE İSTER
y = Cm(8.0)
fc_data(slide, CX, y, "SİSTEM KULLANICIDAN\nSİLME GEREKÇESİ İSTER", w=Cm(6.5))
arrow_down(slide, CX, y + Cm(1.4), y + Cm(2.0))

# 4) GEREKÇEYİ GİRER
y = Cm(10.0)
fc_process(slide, CX, y, "ADMİN SİLME\nGEREKÇESİNİ GİRER")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 5) KAYIT AKTİF LİSTEDEN KALDIRILIR
y = Cm(11.9)
fc_process(slide, CX, y, "KAYIT AKTİF LİSTEDEN\nKALDIRILIR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 6) AUDIT LOG
y = Cm(13.8)
fc_process(slide, CX, y, "AUDIT LOG KAYDI\nOLUŞTURULUR\n(Silinmiş kayıtlar loglanır)")
arrow_down(slide, CX, y + Cm(1.5), y + Cm(2.1))

# 7) BİTİR
y = Cm(15.9)
fc_terminator(slide, CX, y, "BİTİR")

# Açıklama
txBox = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(19), Cm(5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Silme işlemi sadece Admin yetkisine sahip kullanıcılar tarafından yapılabilir. "
    "Admin, sistemdeki kayıtların yanındaki \"Sil\" butonuna tıkladığında, sistem kullanıcıdan "
    "bir silme gerekçesi girmesini ister. Gerekçe girildikten sonra kayıt, aktif listeden "
    "kaldırılır ve denetim izi (Audit Log) kaydı ile birlikte arşivlenir. Böylece silinen kayıtlar "
    "loglama amacıyla sistemde tutulmaya devam eder, ancak aktif listelerde görünmez. "
    "Süreç bu adımla sona erer.")
tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = KOYU_GRI

add_footer(slide, 9)


# ═══════════════════════════════════════════════════════════════
# SLAYT 10: İŞ AKIŞ ŞEMASI – ÇIKIŞ SÜRECİ (ŞEKİLLİ)
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "İŞ AKIŞ ŞEMALARI")
add_section_title(slide, "Çıkış Yapma Süreci Akış Şeması", Cm(2.8))

CX = Cm(10.5)

# 1) BAŞLAT
y = Cm(4.5)
fc_terminator(slide, CX, y, "BAŞLAT")
arrow_down(slide, CX, y + Cm(1.0), y + Cm(1.6))

# 2) ÇIKIŞ YAP BUTONUNA BASILIR
y = Cm(6.1)
fc_process(slide, CX, y, "\"ÇIKIŞ YAP\" BUTONUNA\nBASILIR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 3) ÇIKIŞ SAATİ KAYDEDİLİR
y = Cm(8.0)
fc_process(slide, CX, y, "SİSTEM ÇIKIŞ SAATİNİ\nOTOMATİK KAYDEDER")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 4) İÇERİDE KALIŞ SÜRESİ HESAPLANIR
y = Cm(9.9)
fc_process(slide, CX, y, "İÇERİDE KALIŞ SÜRESİ\nOTOMATİK HESAPLANIR\n(giriş – çıkış farkı)")
arrow_down(slide, CX, y + Cm(1.5), y + Cm(2.1))

# 5) KAYIT GEÇMİŞ KAYITLARA AKTARILIR
y = Cm(12.0)
fc_process(slide, CX, y, "RANDEVU \"TAMAMLANDI\"\nSTATÜSÜNE GEÇİRİLİR")
arrow_down(slide, CX, y + Cm(1.3), y + Cm(1.9))

# 6) AKTİF LİSTEDEN ÇIKARILIR
y = Cm(13.9)
fc_process(slide, CX, y, "KAYIT AKTİF LİSTEDEN\nÇIKARILIR → GEÇMİŞ\nKAYITLARA AKTARILIR")
arrow_down(slide, CX, y + Cm(1.5), y + Cm(2.1))

# 7) BİTİR
y = Cm(16.0)
fc_terminator(slide, CX, y, "BİTİR")

# Açıklama
txBox = slide.shapes.add_textbox(Cm(1), Cm(17.5), Cm(19), Cm(5))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = (
    "Ziyaretini tamamlayan bir kişi için sistemdeki \"Çıkış Yap\" butonuna basılır. "
    "Sistem, çıkış zamanı olarak geçerli tarih ve saati belirler ve bu bilgiyi kayda ekler. "
    "İçeride kalış süresi (giriş saati ile çıkış saati arasındaki fark) otomatik hesaplanır. "
    "Ardından randevu \"Tamamlandı\" statüsüne geçer, aktif ziyaretçi listesinden çıkarılarak "
    "geçmiş kayıtlar arşivine aktarılır. Böylece ziyaretçi çıkışı başarıyla kaydedilmiş olur.")
tf.paragraphs[0].font.size = Pt(9); tf.paragraphs[0].font.name = 'Calibri'
tf.paragraphs[0].font.color.rgb = KOYU_GRI

add_footer(slide, 10)


# ═══════════════════════════════════════
# SLAYT 11: KULLANICI ROLLERİ
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "KULLANICI ROLLERİ VE YETKİLENDİRME")
add_section_title(slide, "Rol Bazlı Erişim Kontrol Tablosu (RBAC)", Cm(3))
add_table(slide, Cm(0.5), Cm(4.2), Cm(20), Cm(11.5), 14, 5,
    [["Özellik / İşlem", "Admin", "Resepsiyon", "Personel", "Sekreter"],
     ["Randevu Oluşturma", "✓", "✓", "✓", "✓ (Yöneticisi adına)"],
     ["Randevu Onaylama / Reddetme", "✓ (Tümü)", "—", "✓ (Kendi)", "—"],
     ["Ziyaretçi Giriş / Çıkış İşlemi", "✓", "✓", "—", "—"],
     ["Tüm Randevuları Görüntüleme", "✓", "✓", "—", "—"],
     ["Kendi Randevularını Görüntüleme", "✓", "✓", "✓", "✓"],
     ["Kara Liste Yönetimi", "✓", "—", "—", "—"],
     ["Kayıt Silme (Gerekçeli)", "✓", "—", "—", "—"],
     ["Audit Log Görüntüleme", "✓", "—", "—", "—"],
     ["Personel Yönetimi", "✓", "—", "—", "—"],
     ["CSV / PDF Dışa Aktarma", "✓", "✓", "—", "—"],
     ["Takvim Yönetimi", "✓", "✓", "✓", "✓"],
     ["Mesajlaşma", "✓", "✓", "✓", "✓"],
     ["Bildirim Alma", "✓", "✓", "✓", "✓"]],
    col_widths=[Cm(6), Cm(3), Cm(3.5), Cm(3.5), Cm(4)])
add_footer(slide, 11)


# ═══════════════════════════════════════
# SLAYT 12: SAYFA ÖN TASARIMLARI 1
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "SAYFA ÖN TASARIMLARI – İŞ KURALLARI")
add_section_title(slide, "1. Giriş (Login) Sayfası", Cm(3))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(4.2), Cm(18), Cm(7))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Giriş Sayfası İş Kuralları:"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI
for rule in [
    "Kullanıcı, kendisine verilen kullanıcı adı ve şifre ile sisteme giriş yapar.",
    "Şifreler veritabanında hash'lenerek (Werkzeug Security – pbkdf2) saklanır.",
    "Belirli sayıda hatalı giriş denemesinde (brute-force koruması) hesap geçici süreyle kilitlenir.",
    "Kilit süresi boyunca kullanıcıya geri sayım bilgisi gösterilir.",
    "Giriş başarılı olduğunda kullanıcının rolüne göre ilgili Dashboard ekranına yönlendirilir."]:
    p = tf.add_paragraph(); p.text = f"• {rule}"
    p.font.size = Pt(9.5); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(4)

add_section_title(slide, "2. Dashboard (Ana Kontrol Paneli)", Cm(12.5))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(13.7), Cm(18), Cm(12))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Dashboard İş Kuralları:"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI
for rule in [
    "Admin: Tüm istatistikleri, logları, kara listeyi ve personel yönetimini görür.",
    "Resepsiyon: Günlük ziyaretçi trafiğini, giriş/çıkış butonlarını ve mesajları görür.",
    "Personel: Kendi takvimini, mesajlarını ve sadece kendi adına oluşturulan randevuları görüp onaylayabilir.",
    "Sekreter: Bağlı olduğu yöneticinin randevularını görür ve adına randevu oluşturabilir.",
    "Yan menüde rol bazlı erişim butonları dinamik olarak gösterilir/gizlenir.",
    "Üst bildirim alanında okunmamış mesajlar ve yeni bildirimler gösterilir."]:
    p = tf.add_paragraph(); p.text = f"• {rule}"
    p.font.size = Pt(9.5); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(4)
add_footer(slide, 12)


# ═══════════════════════════════════════
# SLAYT 13: SAYFA ÖN TASARIMLARI 2
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "SAYFA ÖN TASARIMLARI – İŞ KURALLARI")
add_section_title(slide, "3. Randevu Oluşturma Formu", Cm(3))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(4.2), Cm(18), Cm(7))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Form İş Kuralları:"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI
for rule in [
    "Ziyaretçinin T.C. Kimlik numarası, algoritma ile doğrulanır (11 haneli, matematiksel kontrol).",
    "T.C. Kimlik ve/veya e-posta, kara listede varsa kayıt engellenir ve uyarı gösterilir.",
    "Görüşülecek personel, sistemdeki aktif kullanıcılar listesinden seçilir.",
    "Randevu tarih/saati, tarih-saat seçici (date-time picker) ile belirlenir.",
    "KVKK onay kutusu işaretlenmeden form gönderilemez.",
    "Başarılı kayıt sonrası personele bildirim ve/veya e-posta gönderilir."]:
    p = tf.add_paragraph(); p.text = f"• {rule}"
    p.font.size = Pt(9.5); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(4)

add_section_title(slide, "4. Ziyaretçi Listesi ve Yönetim Sayfaları", Cm(13))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(14.2), Cm(18), Cm(12))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Ziyaretçi Yönetimi İş Kuralları:"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI
for rule in [
    "Randevular; Bekliyor, Onaylandı, Reddedildi, İçeride ve Tamamlandı durumlarıyla takip edilir.",
    "Admin/Resepsiyon, ziyaretçiyi 'İçeride' olarak işaretlediğinde giriş saati otomatik kaydedilir.",
    "Çıkış işleminde çıkış saati kaydedilir ve içeride kalış süresi otomatik hesaplanır.",
    "Tüm kayıtlar; tarih aralığı, durum ve arama terimlerine göre filtrelenebilir.",
    "Toplu veya tekli randevu verileri CSV ve PDF formatında dışa aktarılabilir.",
    "Ziyaretçi detay sayfasından doğrudan e-posta gönderilebilir ve geçmişi görüntülenebilir."]:
    p = tf.add_paragraph(); p.text = f"• {rule}"
    p.font.size = Pt(9.5); p.font.name = 'Calibri'; p.font.color.rgb = KOYU_GRI
    p.space_after = Pt(4)
add_footer(slide, 13)


# ═══════════════════════════════════════
# SLAYT 14: RAPORLAMA
# ═══════════════════════════════════════
slide = add_blank_slide()
add_blue_header_bar(slide, "RAPORLAMA")
add_section_title(slide, "Teknik Rapor ve Dışa Aktarma Yetenekleri", Cm(3))
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(4.5), Cm(18), Cm(22))
tf = txBox.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Ziyaretçi ve Randevu Yönetim Sistemi – Teknik Rapor Özeti"
tf.paragraphs[0].font.size = Pt(13); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].font.color.rgb = KOYU_MAVI

sections = [
    ("1. Genel Bakış",
     "Bu sistem, kurumsal lokasyonlarda ziyaretçi ve randevu süreçlerini dijital ortamda yönetmek amacıyla geliştirilen web tabanlı bir uygulamadır. Python (Flask) backend, Jinja2 template motoru ve SQLite veritabanı kullanılarak oluşturulmuştur. Kullanıcı kimlik doğrulaması, rol bazlı erişim kontrolü, kara liste güvenliği ve denetim izi (audit log) gibi kurumsal güvenlik özellikleri içermektedir."),
    ("2. Kullanıcı Rolleri ve Giriş",
     "Sistem dört tür kullanıcıyı destekler: Admin, Resepsiyon, Personel ve Sekreter. Giriş sonrası arayüz, kullanıcının rolüne göre otomatik olarak şekillenir. Yetkisiz erişim denemeleri engellenir ve loglanır."),
    ("3. Ziyaretçi ve Randevu Kayıt Süreci",
     "Kullanıcılar; ziyaretçinin T.C. Kimlik No, ad-soyad, telefon, e-posta, şirket bilgisi, ziyaret nedeni ve görüşülecek personel bilgilerini girerek randevu oluşturur. T.C. Kimlik doğrulaması ve kara liste kontrolü otomatik yapılır."),
    ("4. Admin ve Resepsiyon Fonksiyonları",
     "Gerekçeyle birlikte kayıt silme, ziyaretçi giriş/çıkış işlemi (check-in/check-out), randevu onay/red, tüm kayıtları görüntüleme, CSV ve PDF dışa aktarma, kara liste ve personel yönetimi."),
    ("5. İletişim ve Bildirim Sistemi",
     "E-posta bildirimleri Flask-Mail (SMTP) üzerinden gönderilir. Platform içi bildirim sistemi ile kullanıcılar anlık olarak bilgilendirilir. Personeller arası mesajlaşma modülü ile kurum içi iletişim sağlanır."),
    ("6. Dışa Aktarma Formatları",
     "CSV: Tüm randevu, ziyaretçi, denetim logu ve kara liste verileri CSV formatında indirilebilir. PDF: ReportLab kütüphanesi ile profesyonel formatlı PDF raporları oluşturulabilir (A4/Landscape)."),
    ("7. Sonuç",
     "Bu sunucu taraflı sistem; rol bazlı erişim kontrolü, T.C. Kimlik doğrulama, kara liste güvenliği, denetim izi, takvim yönetimi, mesajlaşma ve gelişmiş raporlama yetenekleriyle kurumsal düzeyde bir ziyaretçi ve randevu yönetim çözümü sunmaktadır."),
]

for title, desc in sections:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(10); p.font.bold = True
    p.font.name = 'Calibri'; p.font.color.rgb = ORTA_MAVI
    p.space_before = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9); p2.font.name = 'Calibri'; p2.font.color.rgb = KOYU_GRI
    p2.space_after = Pt(4)
add_footer(slide, 14)


# ═══════════════════════════════════════
# KAYDET
# ═══════════════════════════════════════
output_path = "/Users/berkayaktimur/Desktop/ziyaretci_randevu/Analiz_Dokumani.pptx"
prs.save(output_path)
print(f"Analiz dokümanı başarıyla oluşturuldu: {output_path}")
print(f"Toplam slayt sayısı: {len(prs.slides)}")
