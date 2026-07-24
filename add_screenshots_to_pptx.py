from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation('/Users/berkayaktimur/Desktop/ziyaretci_randevu/Analiz_Dokumani.pptx')

screenshots = [
    ('/tmp/screenshots/1_login.png', 'Giriş Ekranı'),
    ('/tmp/screenshots/2_dashboard.png', 'Dashboard (Yönetim Paneli)'),
    ('/tmp/screenshots/3_randevu.png', 'Randevu/Ziyaretçi Oluşturma'),
    ('/tmp/screenshots/4_ziyaretciler.png', 'Ziyaretçiler Listesi'),
    ('/tmp/screenshots/5_personel.png', 'Personel Yönetimi')
]

for img_path, title_text in screenshots:
    if os.path.exists(img_path):
        # 5: Title and Content layout
        slide_layout = prs.slide_layouts[5] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = "Ekran Görüntüsü: " + title_text
        
        # Add image
        # Center the image as best as we can
        try:
            pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
        except Exception as e:
            print(f"Failed to add {img_path}: {e}")

prs.save('/Users/berkayaktimur/Desktop/ziyaretci_randevu/Analiz_Dokumani.pptx')
print("Screenshots added to PPTX.")
