import sys
try:
    from pptx import Presentation
    prs = Presentation(sys.argv[1])
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
except ImportError:
    print("python-pptx not installed")
